from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import json

ppt_path = r"D:\Trae_Solo_Project\拾寻\拾寻 · 城市拾遗网络.pptx"
prs = Presentation(ppt_path)

print(f"Slide width: {prs.slide_width.inches:.2f} inches")
print(f"Slide height: {prs.slide_height.inches:.2f} inches")
print(f"Total slides: {len(prs.slides)}")
print("=" * 80)

def traverse_shapes(shapes, indent=0):
    for shape in shapes:
        prefix = "  " * indent
        if hasattr(shape, "text") and shape.text.strip():
            print(f"{prefix}[TEXT] {shape.text.strip()}")
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print(f"{prefix}[IMAGE] {shape.name}, left={shape.left.inches:.2f}, top={shape.top.inches:.2f}")
        if shape.has_table:
            print(f"{prefix}[TABLE]")
            for row in shape.table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                print(f"{prefix}  |  ".join(cells))
        if shape.has_chart:
            print(f"{prefix}[CHART] {shape.name}")
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            print(f"{prefix}[GROUP] {shape.name}")
            traverse_shapes(shape.shapes, indent + 1)

for idx, slide in enumerate(prs.slides, 1):
    print(f"\n--- Slide {idx} ---")
    traverse_shapes(slide.shapes)
